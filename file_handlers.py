import tempfile
import moviepy.editor as mp
from io import BytesIO
from PIL import Image
import base64
import os
import requests
import streamlit as st
from concurrent.futures import ThreadPoolExecutor, as_completed
from pydub import AudioSegment
from pydub.silence import detect_nonsilent
import time
import docx  # Make sure you have python-docx installed
import pandas as pd
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def convert_video_to_mp3(uploaded_file, suffix):
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_video_file:
            temp_video_file.write(uploaded_file.getbuffer())
            temp_video_file_path = temp_video_file.name

        video = mp.VideoFileClip(temp_video_file_path)

        if video.audio is None:
            st.error("No audio track found in the video file.")
            return None

        with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as audio_file:
            audio_file_path = audio_file.name

        # Using ffmpeg directly for conversion to avoid potential issues
        video.audio.write_audiofile(audio_file_path, codec='mp3')

        return audio_file_path
    except Exception as e:
        st.error(f"Error converting video to audio: {e}")
        return None

def read_docx(file, openai_client):
    doc = docx.Document(file)
    text = "\n".join([para.text for para in doc.paragraphs])
    images = []

    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            image = rel.target_part.blob
            image_stream = BytesIO(image)
            images.append(image_stream)

    image_texts = process_images_concurrently(images, openai_client, "DOCX")
    return text + "\n" + "\n".join(image_texts)

def read_txt(file):
    return file.read().decode("utf-8")

def read_excel(file):
    df = pd.read_excel(file)
    return df.to_string(index=False)

def read_pdf(file, openai_client):
    document = fitz.open(stream=file.read(), filetype="pdf")
    text = ""
    images = []

    for page_num in range(len(document)):
        page = document.load_page(page_num)
        text += page.get_text()
        image_list = page.get_images(full=True)
        for image_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = document.extract_image(xref)
            image_bytes = base_image["image"]
            image_stream = BytesIO(image_bytes)
            images.append(image_stream)

    image_texts = process_images_concurrently(images, openai_client, "PDF")
    return text + "\n" + "\n".join(image_texts)

def read_pptx(file, openai_client):
    presentation = Presentation(file)
    slides = []

    for slide_num, slide in enumerate(presentation.slides, start=1):
        slide_text = f"--- Slide {slide_num} ---\n"
        images = []
        for shape in slide.shapes:
            # Extract text if available
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    slide_text += paragraph.text + "\n"
            
            # Process image if the shape is a picture and has a valid blob
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    if hasattr(shape, 'image') and shape.image.blob:
                        image_stream = BytesIO(shape.image.blob)
                        images.append(image_stream)
                    else:
                        st.warning(f"Slide {slide_num} contains a shape marked as an image but lacks an embedded image.")
                except Exception as e:
                    st.error(f"Error processing an image on slide {slide_num}: {e}")

        # If there are images, transcribe them and append descriptions to the slide text
        if images:
            image_texts = process_images_concurrently(images, openai_client, f"Slide {slide_num}")
            slide_text += "\nImage Descriptions:\n" + "\n".join(image_texts)
        
        slides.append(slide_text)

    return "\n".join(slides)

def process_images_concurrently(images, openai_client, context):
    image_texts = []
    with ThreadPoolExecutor(max_workers=5) as executor:  # Limit to 5 workers
        futures = {executor.submit(transcribe_image, openai_client, image_stream): image_stream for image_stream in images}
        for i, future in enumerate(as_completed(futures)):
            try:
                image_text = future.result()
                image_texts.append(f"Image {i+1}: {image_text}")
            except Exception as e:
                st.error(f"Error processing an image from {context}: {e}")

    return image_texts

def transcribe_image(openai_client, image_stream):
    try:
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {self.api_key}"
        }

        payload = {
            "model": "gpt-4o-mini",
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "What’s in this image?"
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{base64_image}"
                            }
                        }
                    ]
                }
            ],
            "max_tokens": 300
        }

        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload)
        response.raise_for_status()  # This will raise an error for HTTP codes 4xx/5xx

        return response.json()['choices'][0]['message']['content']

    except requests.exceptions.HTTPError as http_err:
        st.error(f"HTTP error occurred: {http_err}")  # Logs the HTTP error in Streamlit
        return f"Error transcribing image: {http_err}"
    except Exception as e:
        st.error(f"An error occurred: {e}")  # Logs any other errors in Streamlit
        return f"Error transcribing image: {e}"

def encode_image(image):
    with BytesIO() as buffer:
        image.save(buffer, format=image.format)
        return base64.b64encode(buffer.getvalue()).decode()


def trim_silence(audio_file, file_name, silence_len=1000, silence_thresh=-40):
    sound = AudioSegment.from_file(audio_file, format="mp3")
    nonsilent_ranges = detect_nonsilent(sound, min_silence_len=silence_len, silence_thresh=silence_thresh)
    
    if nonsilent_ranges:
        start_trim = nonsilent_ranges[0][0]
        end_trim = nonsilent_ranges[-1][1]
        trimmed_sound = sound[start_trim:end_trim]
        
        trimmed_audio_file = BytesIO()
        trimmed_sound.export(trimmed_audio_file, format="mp3")
        trimmed_audio_file.name = file_name  # Set the name attribute for the BytesIO object
        trimmed_audio_file.seek(0)
        return trimmed_audio_file
    return audio_file

def process_files_concurrently(uploaded_files, openai_client):
    transcriptions = []
    with ThreadPoolExecutor() as executor:
        futures = []
        for i, uploaded_file in enumerate(uploaded_files):
            st.info(f"Submitting file {i+1}/{len(uploaded_files)}: {getattr(uploaded_file, 'name', 'unknown')} for processing")
            if uploaded_file.type in ["video/quicktime", "video/mp4"]:
                suffix = ".mov" if uploaded_file.type == "video/quicktime" else ".mp4"
                audio_file_path = convert_video_to_mp3(uploaded_file, suffix)
                if audio_file_path:
                    trimmed_audio_file = trim_silence(audio_file_path, uploaded_file.name)
                    futures.append(executor.submit(openai_client.transcribe_audio, trimmed_audio_file))
            elif uploaded_file.type == "audio/mpeg":
                trimmed_audio_file = trim_silence(uploaded_file, uploaded_file.name)
                futures.append(executor.submit(openai_client.transcribe_audio, trimmed_audio_file))
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                futures.append(executor.submit(read_docx, uploaded_file, openai_client))
            elif uploaded_file.type == "text/plain":
                futures.append(executor.submit(read_txt, uploaded_file))
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                futures.append(executor.submit(read_excel, uploaded_file))
            elif uploaded_file.type == "application/pdf":
                futures.append(executor.submit(read_pdf, uploaded_file, openai_client))
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                futures.append(executor.submit(read_pptx, uploaded_file, openai_client))
            elif uploaded_file.type in ["image/jpeg", "image/png"]:
                base64_image = encode_image(Image.open(uploaded_file))
                futures.append(executor.submit(openai_client.transcribe_image, base64_image))

        for i, future in enumerate(as_completed(futures)):
            try:
                result = future.result()
                transcriptions.append(result)
                st.info(f"Completed processing file {i+1}/{len(uploaded_files)}")
            except Exception as e:
                st.error(f"Error processing file {i+1}/{len(uploaded_files)}: {e}")

    return transcriptions

