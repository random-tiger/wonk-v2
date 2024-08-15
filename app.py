import streamlit as st
import os
from openai import OpenAI
from docx import Document
from io import BytesIO
import moviepy.editor as mp
import tempfile
import docx
import pandas as pd
import fitz
import base64
import requests
from st_aggrid import AgGrid, GridOptionsBuilder, GridUpdateMode
from streamlit_quill import st_quill
from pptx import Presentation
from PIL import Image

api_key = st.secrets["OPENAI_API_KEY"]

# Initialize OpenAI client
client = OpenAI(api_key=api_key)

# Function to transcribe audio using Whisper
def transcribe_audio(audio_file):
    transcription = client.audio.transcriptions.create(model="whisper-1", file=audio_file)
    return transcription['text'] if isinstance(transcription, dict) else transcription.text

# Function to generate response based on prompt and model
def generate_response(transcription, model, custom_prompt):
    response = client.chat.completions.create(
        model=model,
        temperature=0,
        messages=[
            {"role": "system", "content": custom_prompt},
            {"role": "user", "content": transcription}
        ]
    )
    return response.choices[0].message.content

# Function to save meeting minutes as a Word document
def save_as_docx(minutes):
    doc = Document()
    for key, value in minutes.items():
        heading = ' '.join(word.capitalize() for word in key.split('_'))
        doc.add_heading(heading, level=1)
        doc.add_paragraph(value)
        doc.add_paragraph()
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

# Function to convert video files to .mp3
def convert_video_to_mp3(uploaded_file, suffix):
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_video_file:
        temp_video_file.write(uploaded_file.getbuffer())
        temp_video_file_path = temp_video_file.name

    video = mp.VideoFileClip(temp_video_file_path)

    if video.audio is None:
        st.error(f"The uploaded {suffix} file does not contain an audio track.")
        return None

    with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as audio_file:
        audio_file_path = audio_file.name

    video.audio.write_audiofile(audio_file_path)
    return audio_file_path

# Function to read text from a .docx file
def read_docx(file):
    doc = docx.Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

# Function to read text from a .txt file
def read_txt(file):
    return file.read().decode("utf-8")

# Function to read text from an Excel file
def read_excel(file):
    df = pd.read_excel(file)
    return df.to_string(index=False)

# Function to read text from a PDF file
def read_pdf(file):
    document = fitz.open(stream=file.read(), filetype="pdf")
    text = ""
    for page_num in range(len(document)):
        page = document.load_page(page_num)
        text += page.get_text()
    return text

# Function to read text from a PowerPoint file
def read_pptx(file):
    presentation = Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Function to encode image to base64
def encode_image(image):
    with BytesIO() as buffer:
        image.save(buffer, format=image.format)
        return base64.b64encode(buffer.getvalue()).decode()

# Function to transcribe image using GPT-4's multimodal capabilities
def transcribe_image(image_file):
    image = Image.open(image_file)
    base64_image = encode_image(image)

    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {os.getenv('OPENAI_API_KEY')}"
    }

    payload = {
        "model": "gpt-4o-mini",  # Assuming "gpt-4o-mini" is the model with vision capabilities
        "messages": [
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": "What’s in this image?"
                    },
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/jpeg;base64,{base64_image}"
                        }
                    }
                ]
            }
        ],
        "max_tokens": 300
    }

    response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload)
    return response.json()['choices'][0]['message']['content']

# Pre-canned prompts and their respective headings
pre_canned_prompts = {
    "meeting_summary": {
        "summary": {
            "prompt": "You are a highly skilled AI trained in language comprehension and summarization. I would like you to read the following text and summarize it into a concise abstract paragraph. Aim to retain the most important points, providing a coherent and readable summary that could help a person understand the main points of the discussion without needing to read the entire text. Please avoid unnecessary details or tangential points.",
            "heading": "Summary"
        },
        "key_points": {
            "prompt": "You are a proficient AI with a specialty in distilling information into key points. Based on the following text, identify and list the main points that were discussed or brought up. These should be the most important ideas, findings, or topics that are crucial to the essence of the discussion. Your goal is to provide a list that someone could read to quickly understand what was talked about.",
            "heading": "Key Points"
        },
        "action_items": {
            "prompt": "You are an AI expert in analyzing conversations and extracting action items. Please review the text and identify any tasks, assignments, or actions that were agreed upon or mentioned as needing to be done. These could be tasks assigned to specific individuals, or general actions that the group has decided to take. Please list these action items clearly and concisely.",
            "heading": "Action Items"
        },
        "sentiment": {
            "prompt": "As an AI with expertise in language and emotion analysis, your task is to analyze the sentiment of the following text. Please consider the overall tone of the discussion, the emotion conveyed by the language used, and the context in which words and phrases are used. Indicate whether the sentiment is generally positive, negative, or neutral, and provide brief explanations for your analysis where possible.",
            "heading": "Sentiment Analysis"
        }
    },
    "user_research": {
        "summary": {
            "prompt": "You are a highly skilled AI trained in language comprehension and summarization. I would like you to read the following text and summarize it into a concise abstract paragraph. Aim to retain the most important points, providing a coherent and readable summary that could help a person understand the main points of the discussion without needing to read the entire text. Please avoid unnecessary details or tangential points.",
            "heading": "Summary"
        },
        "biographical_info": {
            "prompt": "You are a proficient AI with a specialty in distilling biographical information about people. Based on the following text, please identify biographical information about the subject of the research study.",
            "heading": "Biographical Info"
        },
        "key_insights": {
            "prompt": "You are a proficient AI with a specialty in distilling information into key points. Based on the following user research transcript, please identify the key insights. Identify and list the main points that were discussed or brought up. These should be the most important ideas, findings, or topics that are crucial to the essence of the discussion. Your goal is to provide a list that someone could read to quickly understand what was talked about.",
            "heading": "Key Insights"
        },
        "recommendations": {
            "prompt": "You are a proficient AI with a specialty in identifying meaningful product opportunities. Based on the transcript, please identify product recommendations/opportunities.",
            "heading": "Recommendations"
        }
    },
    "action_items": {
        "action_items": {
            "prompt": "You are an AI expert in analyzing conversations and extracting action items. Please review the text and identify any tasks, assignments, or actions that were agreed upon or mentioned as needing to be done. These could be tasks assigned to specific individuals, or general actions that the group has decided to take. Please list these action items clearly and concisely. Each task should be formatted as follows: Topic. Description of the task. Do not use sub-bullets.",
            "heading": "Action Items"
        }
    }
}

# Streamlit app
def main():
    # Custom CSS to set app width and reduce gutter space
    st.markdown(
        """
        <style>
        .reportview-container .main .block-container {
            padding-left: 0rem;
            padding-right: 0rem;
            max-width: 100%;
            margin: 0 auto;
        }
        .css-18e3th9 {
            flex: 1 1 100%;
            width: 100%;
            padding: 2rem 1rem 1rem;
        }
        </style>
        """,
        unsafe_allow_html=True
    )

    st.sidebar.title("Wonk")

    st.sidebar.info("Upload mp3, mp4, mov, docx, txt, xlsx, pdf, pptx, or image files to start!")
    uploaded_files = st.sidebar.file_uploader("Upload audio, video, text, or image files", type=["mp3", "mp4", "mov", "docx", "txt", "xlsx", "pdf", "pptx", "jpg", "jpeg", "png"], accept_multiple_files=True)
    process_files = st.sidebar.button("Process Files")

    if uploaded_files is not None and process_files:
        if "transcriptions" not in st.session_state:
            st.session_state.transcriptions = []

        for uploaded_file in uploaded_files:
            if uploaded_file.type in ["video/quicktime", "video/mp4"]:
                suffix = ".mov" if uploaded_file.type == "video/quicktime" else ".mp4"
                audio_file_path = convert_video_to_mp3(uploaded_file, suffix)
                if audio_file_path is not None:
                    with open(audio_file_path, "rb") as f:
                        st.session_state.transcriptions.append(transcribe_audio(f))
            elif uploaded_file.type == "audio/mpeg":
                st.session_state.transcriptions.append(transcribe_audio(uploaded_file))
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                st.session_state.transcriptions.append(read_docx(uploaded_file))
            elif uploaded_file.type == "text/plain":
                st.session_state.transcriptions.append(read_txt(uploaded_file))
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                st.session_state.transcriptions.append(read_excel(uploaded_file))
            elif uploaded_file.type == "application/pdf":
                st.session_state.transcriptions.append(read_pdf(uploaded_file))
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                st.session_state.transcriptions.append(read_pptx(uploaded_file))
            elif uploaded_file.type in ["image/jpeg", "image/png"]:
                st.session_state.transcriptions.append(transcribe_image(uploaded_file))

        if st.session_state.transcriptions:
            combined_transcription = "\n\n".join(st.session_state.transcriptions)
            st.session_state.transcription = combined_transcription

    if "transcription" in st.session_state:
        transcription = st.session_state.transcription
        with st.expander("Transcription", expanded=True):
            st.subheader("Transcription")
            edited_transcription = st_quill(value=transcription, key='transcription_editor')
            st.session_state.transcription = edited_transcription

        st.sidebar.info("Select what you'd like to create!")
        summary_type = st.sidebar.radio(
            "Select the type of summary you want to generate:",
            ("", "Meeting Summary", "User Research Synthesis", "Action Items"),
            index=0
        )

        if 'prompts' not in st.session_state:
            st.session_state.prompts = []

        checkboxes = {}
        if summary_type == "Meeting Summary":
            st.sidebar.markdown("### Meeting Summary Prompts")
            st.sidebar.info("Select the sections you'd like in your document!")
            checkboxes = {
                "summary": st.sidebar.checkbox("Summary"),
                "key_points": st.sidebar.checkbox("Key Points"),
                "action_items": st.sidebar.checkbox("Action Items"),
                "sentiment": st.sidebar.checkbox("Sentiment Analysis")
            }

        elif summary_type == "User Research Synthesis":
            st.sidebar.markdown("### User Research Synthesis Prompts")
            st.sidebar.info("Select the sections you'd like in your document!")
            checkboxes = {
                "summary": st.sidebar.checkbox("Summary", key="user_summary"),
                "biographical_info": st.sidebar.checkbox("Biographical Info"),
                "key_insights": st.sidebar.checkbox("Key Insights"),
                "recommendations": st.sidebar.checkbox("Recommendations")
            }

        elif summary_type == "Action Items":
            st.sidebar.markdown("### Action Items Prompt")
            st.sidebar.info("Select the section to generate action items!")
            checkboxes = {
                "action_items": st.sidebar.checkbox("Action Items", key="action_items")
            }

        if any(checkboxes.values()):
            st.sidebar.info("Click 'Create GPT Tasks' to proceed")
            if st.sidebar.button("Create GPT Tasks"):
                for key, checked in checkboxes.items():
                    if checked:
                        try:
                            st.session_state.prompts.append({
                                "prompt": pre_canned_prompts[summary_type.lower().replace(" ", "_")][key]["prompt"],
                                "model": "gpt-4o",
                                "heading": pre_canned_prompts[summary_type.lower().replace(" ", "_")][key]["heading"]
                            })
                        except KeyError as e:
                            st.error(f"KeyError: {e} - summary_type: {summary_type.lower().replace(' ', '_')}, key: {key}")
                            st.stop()

        for i, prompt_info in enumerate(st.session_state.prompts):
            with st.expander(f"GPT Task {i+1} - {prompt_info['heading']}", expanded=True):
                st.info("Update the pre-canned prompt to customize!")
                prompt_info["model"] = st.text_input("Model", value=prompt_info["model"], key=f"model_{i}")
                prompt_info["prompt"] = st.text_area("Prompt", value=prompt_info["prompt"], key=f"prompt_{i}")
                if st.button("Remove GPT Task", key=f"remove_gpt_task_{i}"):
                    st.session_state.prompts.pop(i)
                    break

        if st.session_state.prompts:
            st.info("Click generate to create your document!")
            st.markdown(
                """
                <style>
                .blue-button button {
                    background-color: #007BFF !important;
                    color: white !important;
                }
                </style>
                """,
                unsafe_allow_html=True
            )
            if st.button("Generate", key="generate", help=None, on_click=None, disabled=False, use_container_width=False):
                minutes = {}
                for i, prompt_info in enumerate(st.session_state.prompts):
                    task_key = prompt_info["heading"] if prompt_info["heading"] else f"Task {i+1}"
                    minutes[task_key] = generate_response(st.session_state.transcription, prompt_info["model"], prompt_info["prompt"])
                st.session_state.generated_minutes = minutes  # Store the generated minutes in session state

        # Display generated minutes if they exist in session state
        if 'generated_minutes' in st.session_state:
            with st.expander("Generated Minutes", expanded=True):
                for key, value in st.session_state.generated_minutes.items():
                    st.write(f"**{key}**")
                    st.write(value)

                docx_file = save_as_docx(st.session_state.generated_minutes)

                st.info("Click download to get a docx file of your document!")
                st.download_button(
                    label="Download Meeting Minutes",
                    data=docx_file,
                    file_name="meeting_minutes.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                )

            if "Action Items" in st.session_state.generated_minutes:
                with st.expander("Action Items", expanded=True):
                    st.subheader("Action Items")
                    action_items = st.session_state.generated_minutes["Action Items"]
                    st.info("Check boxes to generate documents from tasks!")
                    action_items_list = action_items.split('\n')
                    action_items_list = [item for item in action_items_list if item]  # Remove empty items

                    action_items_dict = {}
                    parent_task = None

                    for item in action_items_list:
                        if item.startswith("    "):  # Child task
                            if parent_task:
                                action_items_dict[parent_task].append(item.strip())
                        else:  # Parent task
                            parent_task = item.strip()
                            action_items_dict[parent_task] = []

                    # Prepare data for AgGrid
                    grid_data = []
                    for idx, (parent, children) in enumerate(action_items_dict.items(), 1):
                        grid_data.append({
                            "Task Number": idx,
                            "Task": parent,
                            "Draft Email": False,
                            "Draft Slack": False,
                            "Draft Memo": False
                        })

                    # Convert list of dicts to DataFrame
                    grid_df = pd.DataFrame(grid_data)

                    # Configure AgGrid
                    gb = GridOptionsBuilder.from_dataframe(grid_df)
                    gb.configure_column("Draft Email", editable=True, cellEditor="agCheckboxCellEditor")
                    gb.configure_column("Draft Slack", editable=True, cellEditor="agCheckboxCellEditor")
                    gb.configure_column("Draft Memo", editable=True, cellEditor="agCheckboxCellEditor")
                    gb.configure_pagination()
                    gb.configure_default_column(editable=True, resizable=True)
                    grid_options = gb.build()

                    grid_response = AgGrid(grid_df, gridOptions=grid_options, height=300, fit_columns_on_grid_load=True, update_mode=GridUpdateMode.MODEL_CHANGED)

                    # Handling the checkbox responses
                    if isinstance(grid_response['data'], pd.DataFrame):
                        for index, row in grid_response['data'].iterrows():
                            if row["Draft Email"]:
                                st.session_state[f"email_prompt_{row['Task Number']}"] = f"Draft an email for the following action item: {row['Task']}"
                                row["Draft Email"] = False
                            if row["Draft Slack"]:
                                st.session_state[f"slack_prompt_{row['Task Number']}"] = f"Draft a Slack message for the following action item: {row['Task']}"
                                row["Draft Slack"] = False
                            if row["Draft Memo"]:
                                st.session_state[f"memo_prompt_{row['Task Number']}"] = f"Draft a memo for the following action item: {row['Task']}"
                                row["Draft Memo"] = False

                    # Display generated drafts
                    #st.write("### Generate Drafts")
                    for key in st.session_state.keys():
                        if key.startswith("email_prompt_"):
                            task_num = key.split('_')[-1]
                            st.subheader(f"Email Draft for Task {task_num}")
                            st.write(st.session_state[key])
                            if st.button(f"Generate Email for Task {task_num}"):
                                draft = generate_response(st.session_state.transcription, "gpt-4o", st.session_state[key])
                                st.write(draft)
                        elif key.startswith("slack_prompt_"):
                            task_num = key.split('_')[-1]
                            st.subheader(f"Slack Draft for Task {task_num}")
                            st.write(st.session_state[key])
                            if st.button(f"Generate Slack for Task {task_num}"):
                                draft = generate_response(st.session_state.transcription, "gpt-4o", st.session_state[key])
                                st.write(draft)
                        elif key.startswith("memo_prompt_"):
                            task_num = key.split('_')[-1]
                            st.subheader(f"Memo Draft for Task {task_num}")
                            st.write(st.session_state[key])
                            if st.button(f"Generate Memo for Task {task_num}"):
                                draft = generate_response(st.session_state.transcription, "gpt-4o", st.session_state[key])
                                st.write(draft)

if __name__ == "__main__":
    main()
